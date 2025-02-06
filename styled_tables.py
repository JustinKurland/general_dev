import pandas as pd
from html2image import Html2Image
from pptx import Presentation
from pptx.util import Inches

def save_styled_dataframe_as_image(styled_df, filename="styled_table.png"):
    """Converts a styled Pandas DataFrame to an image while preserving styles."""
    
    # Convert styled DataFrame to HTML
    html = styled_df.render()

    # Save the HTML to a temporary file
    html_filename = "styled_table.html"
    with open(html_filename, "w", encoding="utf-8") as f:
        f.write(html)

    # Convert the HTML to an image
    hti = Html2Image(output_path=".")
    hti.screenshot(html_file=html_filename, save_as=filename)

def add_image_to_ppt(image_path, ppt_filename="output.pptx"):
    """Inserts an image into a PowerPoint slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    left, top = Inches(1), Inches(1)  # Positioning
    slide.shapes.add_picture(image_path, left, top, width=Inches(6))
    prs.save(ppt_filename)

# Example usage: Apply styling to a DataFrame
df = pd.DataFrame({
    "Column A": ["Row 1", "Row 2", "Row 3"],
    "Column B": [10, 20, 30],
    "Column C": ["✅", "❌", "✅"]
})

# Define styling
styled_df = df.style.set_properties(**{
    "background-color": "lightblue",
    "color": "black",
    "border": "1px solid black",
    "text-align": "center"
})

# Save styled DataFrame as an image
save_styled_dataframe_as_image(styled_df, "styled_table.png")

# Insert the image into PowerPoint
add_image_to_ppt("styled_table.png", "output.pptx")
