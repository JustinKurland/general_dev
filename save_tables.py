import pandas as pd
import matplotlib.pyplot as plt
from pandas.plotting import table
from pptx import Presentation
from pptx.util import Inches

def save_dataframe_as_image(df, filename="table.png"):
    """Converts a Pandas DataFrame to an image and saves it."""
    fig, ax = plt.subplots(figsize=(df.shape[1], df.shape[0]))  # Adjust size dynamically
    ax.set_frame_on(False)  # Remove axes borders
    ax.xaxis.set_visible(False)
    ax.yaxis.set_visible(False)
    
    table_obj = table(ax, df, loc="center", cellLoc="center", colWidths=[0.2] * len(df.columns))
    table_obj.auto_set_font_size(False)
    table_obj.set_fontsize(10)
    
    plt.savefig(filename, bbox_inches="tight", dpi=300)
    plt.close(fig)

def add_image_to_ppt(image_path, ppt_filename="output.pptx"):
    """Inserts an image into a PowerPoint slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide
    left, top = Inches(1), Inches(1)  # Positioning
    slide.shapes.add_picture(image_path, left, top, width=Inches(6))
    prs.save(ppt_filename)

# Example usage
df = pd.DataFrame({
    "Column A": ["Row 1", "Row 2", "Row 3"],
    "Column B": [10, 20, 30],
    "Column C": ["✅", "❌", "✅"]
})

# Save as image and insert into PowerPoint
save_dataframe_as_image(df, "table.png")
add_image_to_ppt("table.png", "output.pptx")



import pandas as pd
import matplotlib.pyplot as plt
from pandas.plotting import table
from pptx import Presentation
from pptx.util import Inches

def save_dataframe_as_image_with_title(df, title="Table Title", filename="table.png"):
    """Converts a Pandas DataFrame to an image with a left-aligned title."""
    
    fig, ax = plt.subplots(figsize=(df.shape[1] + 1, df.shape[0] + 1))  # Adjust size
    ax.set_frame_on(False)  # Remove axes borders
    ax.xaxis.set_visible(False)
    ax.yaxis.set_visible(False)

    # Add table
    table_obj = table(ax, df, loc="center", cellLoc="center", colWidths=[0.2] * len(df.columns))
    table_obj.auto_set_font_size(False)
    table_obj.set_fontsize(10)

    # Add title (left-aligned)
    plt.text(
        x=0, y=1.2, s=title, fontsize=14, fontweight="bold", ha="left", transform=ax.transAxes
    )

    plt.savefig(filename, bbox_inches="tight", dpi=300)
    plt.close(fig)

def add_image_to_ppt(image_path, ppt_filename="output.pptx"):
    """Inserts an image into a PowerPoint slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    left, top = Inches(1), Inches(1)  # Positioning
    slide.shapes.add_picture(image_path, left, top, width=Inches(6))
    prs.save(ppt_filename)

# Example usage
df = pd.DataFrame({
    "Column A": ["Row 1", "Row 2", "Row 3"],
    "Column B": [10, 20, 30],
    "Column C": ["✅", "❌", "✅"]
})

# Save table with title and insert into PowerPoint
save_dataframe_as_image_with_title(df, title="My Left-Aligned Table Title", filename="table_with_title.png")
add_image_to_ppt("table_with_title.png", "output_with_title.pptx")
