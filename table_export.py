import pandas as pd
import xlsxwriter

# Create a sample styled DataFrame
df = pd.DataFrame({
    "Risk Category": ["Low", "Medium", "High"],
    "Count": [10, 20, 30]
})

# Apply styling
styled_df = df.style.set_caption("Incident Frequency Report").set_properties(**{
    "background-color": "lightblue",
    "color": "black",
    "border": "1px solid black",
    "text-align": "center"
})

# Export to Excel with proper styling
with pd.ExcelWriter("output.xlsx", engine="xlsxwriter") as writer:
    workbook = writer.book
    worksheet = workbook.add_worksheet("Risk Level")
    
    # Define Formats
    caption_format = workbook.add_format({'bold': True, 'font_size': 14, 'align': 'left'})
    header_format = workbook.add_format({'bold': True, 'font_size': 12, 'bg_color': '#D9EAD3', 'border': 0})
    cell_format = workbook.add_format({'font_size': 11, 'border': 0})
    
    start_row = 0
    for i in range(1, 13):  # Loop through multiple tables
        styled_df = globals()[f"styled_df_transposed{i}"]

        # Write caption manually
        worksheet.write(start_row, 0, f"Incident Frequency Report {i}", caption_format)
        start_row += 1  # Leave space for the title

        # Write DataFrame to Excel
        styled_df.data.to_excel(writer, sheet_name="Risk Level", startrow=start_row, startcol=0, index=True, header=True)
        start_row += styled_df.data.shape[0] + 3  # Add spacing between tables

        # Format headers and cells
        for col_num, col_name in enumerate(styled_df.data.columns):
            worksheet.write(start_row - styled_df.data.shape[0] - 2, col_num + 1, col_name, header_format)  # Header formatting
            max_len = max(styled_df.data[col_name].astype(str).map(len).max(), len(str(col_name))) + 2
            worksheet.set_column(col_num + 1, col_num + 1, max_len, cell_format)  # Auto-adjust column width

# ✅ Fixes:
# - Titles (captions) now appear above tables
# - Fonts & colors are properly applied
# - Borders around headers removed
# - Tables spaced out correctly

import dataframe_image as dfi
from pptx import Presentation
from pptx.util import Inches

# Function to save styled DataFrame as an image
def save_styled_dataframe_as_image(styled_df, filename="table.png"):
    dfi.export(styled_df, filename)

# Function to insert images into PowerPoint
def add_images_to_ppt(images, ppt_filename="output.pptx"):
    prs = Presentation()
    for img in images:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
        left, top = Inches(1), Inches(1)  # Positioning
        slide.shapes.add_picture(img, left, top, width=Inches(6))
    prs.save(ppt_filename)

# Example usage
images = []
for i in range(1, 13):
    filename = f"table_{i}.png"
    save_styled_dataframe_as_image(globals()[f"styled_df_transposed{i}"], filename)
    images.append(filename)

add_images_to_ppt(images, "styled_presentation.pptx")
